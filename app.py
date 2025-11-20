from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pptx import Presentation
from typing import Dict, Any, List
import tempfile
import os
import re
import subprocess
import json
import shutil
from pathlib import Path

app = FastAPI(title="AIStoryteller Backend")

# CORS配置
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ========== 配置 ==========
SCRIPTS_DIR = Path(__file__).parent / "pptx_skills" / "scripts"
TEMPLATE_PATH = Path(__file__).parent / "templates" / "TeamsPPTTemplate.pptx"
TEMP_DIR = Path(__file__).parent / "temp"
FIXED_TEMPLATE_PATH = Path(__file__).parent / "fixed_template.pptx"  # 新增固定模板路径

# 确保临时目录存在
TEMP_DIR.mkdir(exist_ok=True)

# ========== 基础工具函数 ==========

def shape_type_to_str(shape_type) -> str:
    try:
        return shape_type.name
    except Exception:
        return str(shape_type)


def safe_int(value, default=0) -> int:
    try:
        return int(value)
    except Exception:
        return default


# ========== Slide 分类规则 ==========

AGENDA_KEYWORDS = ["agenda", "contents", "outline", "目录", "议程"]
ENDING_KEYWORDS = ["thank you", "thanks", "q&a", "questions", "谢谢", "结束"]
SECTION_KEYWORDS = ["section", "part", "chapter", "模块", "篇", "part "]


def is_bullet_line(line: str) -> bool:
    line = line.strip()
    if not line:
        return False
    if line.startswith(("•", "-", "·", "· ", "●", "○", "▶")):
        return True
    if re.match(r"^(\d+|[a-zA-Z])[\.\)]\s+", line):
        return True
    return False


def classify_slide(slide_dict: Dict[str, Any], meta: Dict[str, Any]) -> str:
    """分类幻灯片类型"""
    layout_name = (slide_dict.get("layout_name") or "").lower()
    shapes = slide_dict.get("shapes", [])
    slide_h = meta.get("slide_height_emu") or 1
    slide_w = meta.get("slide_width_emu") or 1

    text_shapes: List[Dict[str, Any]] = []
    picture_shapes: List[Dict[str, Any]] = []

    total_text_len = 0
    total_lines = 0
    bullet_lines = 0

    for s in shapes:
        if s.get("shape_type", "").upper() in ("PICTURE", "MEDIA"):
            picture_shapes.append(s)

        if s.get("has_text_frame") and s.get("text"):
            txt = s["text"].strip()
            if not txt:
                continue
            text_shapes.append(s)
            total_text_len += len(txt)

            lines = [ln for ln in txt.splitlines() if ln.strip()]
            total_lines += len(lines)
            bullet_lines += sum(1 for ln in lines if is_bullet_line(ln))

    bullet_ratio = (bullet_lines / total_lines) if total_lines > 0 else 0.0

    # 提取候选大标题
    title_candidate = None
    title_score = -1.0
    for s in text_shapes:
        geom = s.get("geometry", {})
        top = geom.get("top_emu") or 0
        height = geom.get("height_emu") or 0
        text = (s.get("text") or "").strip()

        y_center_ratio = (top + height / 2) / slide_h
        area = (geom.get("width_emu") or 0) * height

        score = area
        if y_center_ratio < 0.35:
            score *= 1.3

        if score > title_score and 0 < len(text) <= 60:
            title_score = score
            title_candidate = text

    all_text = " ".join([(s.get("text") or "") for s in text_shapes]).lower()

    # 分类规则
    if "title" in layout_name and "agenda" not in layout_name:
        return "title"
    if "agenda" in layout_name or "目录" in layout_name:
        return "agenda"

    if slide_dict.get("index") == 0:
        if title_candidate and len(text_shapes) <= 3:
            return "title"

    if len(text_shapes) <= 2 and total_text_len <= 80 and title_candidate:
        return "title"

    if any(k in all_text for k in AGENDA_KEYWORDS):
        return "agenda"
    if bullet_ratio >= 0.5 and any(k in all_text for k in ["agenda", "目录", "contents"]):
        return "agenda"

    if len(text_shapes) <= 2 and total_text_len <= 60:
        if any(k in (title_candidate or "").lower() for k in SECTION_KEYWORDS):
            return "section"

    if any(k in all_text for k in ENDING_KEYWORDS):
        return "ending"

    if picture_shapes and total_text_len <= 120:
        return "content_image"

    if bullet_ratio >= 0.4:
        return "content_bullets"

    if total_text_len > 0:
        return "content"

    return "other"


def extract_ppt_structure(prs: Presentation) -> Dict[str, Any]:
    """提取 PPT 结构"""
    meta = {
        "slide_count": len(prs.slides),
        "slide_width_emu": safe_int(prs.slide_width),
        "slide_height_emu": safe_int(prs.slide_height),
    }

    slides_info: List[Dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_dict: Dict[str, Any] = {
            "index": slide_idx,
            "layout_name": slide.slide_layout.name,
            "shapes": [],
        }

        for shape_idx, shape in enumerate(slide.shapes):
            geom = {
                "left_emu": safe_int(getattr(shape, "left", 0)),
                "top_emu": safe_int(getattr(shape, "top", 0)),
                "width_emu": safe_int(getattr(shape, "width", 0)),
                "height_emu": safe_int(getattr(shape, "height", 0)),
            }

            shape_dict: Dict[str, Any] = {
                "index": shape_idx,
                "name": getattr(shape, "name", None),
                "shape_type": shape_type_to_str(getattr(shape, "shape_type", None)),
                "geometry": geom,
                "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
                "text": None,
            }

            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                shape_dict["text"] = text or None

            slide_dict["shapes"].append(shape_dict)

        slide_dict["slide_type"] = classify_slide(slide_dict, meta)
        slides_info.append(slide_dict)

    return {
        "meta": meta,
        "slides": slides_info,
    }


# ========== 模板匹配 ==========

def match_template_slide(slide_type: str, slide_index: int, total_slides: int) -> int:
    """
    根据用户幻灯片类型智能匹配模板页面索引
    
    返回：模板的幻灯片索引（0-32）
    """
    # 模板映射规则
    template_map = {
        "title": [0, 1, 4],  # 封面：标准、渐变、纯蓝
        "agenda": [6, 7],    # 议程页
        "section": [15, 16, 17],  # 分节页：纯蓝、渐变、白色
        "content_bullets": [18, 19, 20],  # 内容页
        "content_image": [21, 22],  # 图片页
        "content": [18, 19, 20],  # 普通内容
        "ending": [31, 32],  # 结束页
        "other": [18],  # 默认内容页
    }
    
    # 获取候选模板
    candidates = template_map.get(slide_type, [18])
    
    # 策略：第一页用第一个候选，其他循环使用
    if slide_index == 0 and slide_type == "title":
        return candidates[0]
    
    # 循环选择
    idx = slide_index % len(candidates)
    return candidates[idx]


def generate_replacement_json(user_slides: List[Dict], template_inventory: Dict) -> Dict:
    """
    生成替换 JSON
    
    Args:
        user_slides: 用户 PPT 的幻灯片数据
        template_inventory: 模板的 inventory 数据
        
    Returns:
        replacement JSON 格式的字典
    """
    replacement = {}
    
    for slide_idx, user_slide in enumerate(user_slides):
        slide_key = f"slide-{slide_idx}"
        replacement[slide_key] = {}
        
        # 获取该页的模板 shapes
        if slide_key not in template_inventory:
            continue
            
        template_shapes = template_inventory[slide_key]
        
        # 提取用户文本
        user_texts = []
        for shape in user_slide.get("shapes", []):
            if shape.get("has_text_frame") and shape.get("text"):
                text = shape["text"].strip()
                if text:
                    user_texts.append(text)
        
        if not user_texts:
            continue
        
        # 分配文本到模板 shapes
        # 策略：第一个文本作为标题，其他作为内容
        shape_keys = sorted(template_shapes.keys(), 
                          key=lambda k: int(k.split("-")[1]))
        
        for i, shape_key in enumerate(shape_keys):
            if i >= len(user_texts):
                break
                
            text = user_texts[i]
            template_shape = template_shapes[shape_key]
            
            # 判断是标题还是内容
            is_title = (i == 0 or 
                       template_shape.get("placeholder_type") in ["TITLE", "CENTER_TITLE"])
            
            # 判断是否为列表
            lines = [ln.strip() for ln in text.split('\n') if ln.strip()]
            is_bullet_list = len(lines) > 1 or any(is_bullet_line(ln) for ln in lines)
            
            paragraphs = []
            
            if is_bullet_list and not is_title:
                # 列表内容
                for line in lines:
                    # 移除手动的 bullet 符号
                    clean_line = re.sub(r'^[•\-●○▶]\s*', '', line)
                    clean_line = re.sub(r'^\d+[\.\)]\s*', '', clean_line)
                    clean_line = re.sub(r'^[a-zA-Z][\.\)]\s*', '', clean_line)
                    
                    paragraphs.append({
                        "text": clean_line,
                        "bullet": True,
                        "level": 0
                    })
            else:
                # 标题或普通文本
                para = {"text": text}
                if is_title:
                    para["bold"] = True
                    if template_shape.get("placeholder_type") == "CENTER_TITLE":
                        para["alignment"] = "CENTER"
                paragraphs.append(para)
            
            replacement[slide_key][shape_key] = {
                "paragraphs": paragraphs
            }
    
    return replacement


# ========== API 路由 ==========

@app.get("/health")
def health_check():
    return {"status": "ok"}


@app.post("/api/parse_ppt")
async def parse_ppt(file: UploadFile = File(...)):
    """解析用户 PPTX"""
    suffix = ".pptx"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        contents = await file.read()
        tmp.write(contents)
        tmp_path = tmp.name

    try:
        prs = Presentation(tmp_path)
        data = extract_ppt_structure(prs)
        return data
    finally:
        os.remove(tmp_path)


@app.post("/api/beautify_ppt")
async def beautify_ppt(file: UploadFile = File(...)):
    """
    美化用户 PPTX - 修改为始终返回固定模板
    """
    # 检查固定模板文件是否存在
    if not FIXED_TEMPLATE_PATH.exists():
        raise HTTPException(status_code=500, detail="Fixed template file not found")
    
    # 直接返回固定模板文件，忽略用户上传的文件
    return FileResponse(
        path=str(FIXED_TEMPLATE_PATH),
        filename=f"beautified_{file.filename}",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


# ========== 新增的固定模板接口 ==========

@app.get("/api/fixed_template_data")
async def get_fixed_template_data():
    """获取固定模板的解析数据"""
    if not FIXED_TEMPLATE_PATH.exists():
        raise HTTPException(status_code=500, detail="Fixed template file not found")
    
    # 解析固定模板
    prs = Presentation(str(FIXED_TEMPLATE_PATH))
    data = extract_ppt_structure(prs)
    
    return data


@app.get("/fixed_template.pptx")
async def get_fixed_template():
    """提供固定模板文件下载"""
    if not FIXED_TEMPLATE_PATH.exists():
        raise HTTPException(status_code=404, detail="Fixed template not found")
    
    return FileResponse(
        path=str(FIXED_TEMPLATE_PATH),
        filename="beautified_presentation.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)