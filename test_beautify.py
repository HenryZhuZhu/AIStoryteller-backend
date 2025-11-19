#!/usr/bin/env python3
"""
测试 PPT 美化功能

用法:
    python test_beautify.py input.pptx
"""

import sys
import subprocess
import json
from pathlib import Path

# 添加当前目录到 Python 路径
sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation

def test_beautify(input_pptx):
    """测试完整的美化流程"""
    
    print("=" * 60)
    print("PPT 美化功能测试")
    print("=" * 60)
    
    # 配置路径
    SCRIPTS_DIR = Path("pptx_skills/scripts")
    TEMPLATE_PATH = Path("templates/TeamsPPTTemplate.pptx")
    TEMP_DIR = Path("temp")
    
    TEMP_DIR.mkdir(exist_ok=True)
    
    input_path = Path(input_pptx)
    if not input_path.exists():
        print(f"❌ 输入文件不存在: {input_pptx}")
        return
    
    if not TEMPLATE_PATH.exists():
        print(f"❌ 模板文件不存在: {TEMPLATE_PATH}")
        return
    
    print(f"\n📄 输入文件: {input_path}")
    print(f"📐 模板文件: {TEMPLATE_PATH}")
    
    # 步骤 1: 解析用户 PPT
    print("\n[1/6] 解析用户 PPT...")
    prs = Presentation(str(input_path))
    print(f"   ✓ 共 {len(prs.slides)} 页")
    
    # 步骤 2: 模拟分类（简化版）
    print("\n[2/6] 分析页面类型...")
    template_sequence = []
    for i, slide in enumerate(prs.slides):
        # 简单策略：第一页用封面，其他用内容页
        if i == 0:
            template_idx = 0  # 封面
        else:
            template_idx = 18  # 内容页
        template_sequence.append(template_idx)
        print(f"   Slide {i} -> Template {template_idx}")
    
    # 步骤 3: 重排模板
    print("\n[3/6] 重排模板页面...")
    working_pptx = TEMP_DIR / "working.pptx"
    sequence_str = ",".join(map(str, template_sequence))
    
    cmd = [
        "python",
        str(SCRIPTS_DIR / "rearrange.py"),
        str(TEMPLATE_PATH),
        str(working_pptx),
        sequence_str
    ]
    
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"   ❌ 失败: {result.stderr}")
        return
    print(f"   ✓ 完成")
    
    # 步骤 4: 提取 inventory
    print("\n[4/6] 提取模板结构...")
    inventory_json = TEMP_DIR / "inventory.json"
    
    cmd = [
        "python",
        str(SCRIPTS_DIR / "inventory.py"),
        str(working_pptx),
        str(inventory_json)
    ]
    
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"   ❌ 失败: {result.stderr}")
        return
    
    with open(inventory_json, "r", encoding="utf-8") as f:
        inventory = json.load(f)
    print(f"   ✓ 找到 {len(inventory)} 页，{sum(len(v) for v in inventory.values())} 个形状")
    
    # 步骤 5: 生成替换 JSON（简化版）
    print("\n[5/6] 生成替换内容...")
    replacement = {}
    
    for i, slide in enumerate(prs.slides):
        slide_key = f"slide-{i}"
        if slide_key not in inventory:
            continue
        
        replacement[slide_key] = {}
        
        # 提取第一个文本框的内容
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # 假设第一个文本框是标题
                replacement[slide_key]["shape-0"] = {
                    "paragraphs": [
                        {
                            "text": shape.text.strip(),
                            "bold": True
                        }
                    ]
                }
                break
    
    replacement_json = TEMP_DIR / "replacement.json"
    with open(replacement_json, "w", encoding="utf-8") as f:
        json.dump(replacement, f, ensure_ascii=False, indent=2)
    print(f"   ✓ 生成替换内容")
    
    # 步骤 6: 应用替换
    print("\n[6/6] 应用替换...")
    output_pptx = TEMP_DIR / "output.pptx"
    
    cmd = [
        "python",
        str(SCRIPTS_DIR / "replace.py"),
        str(working_pptx),
        str(replacement_json),
        str(output_pptx)
    ]
    
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"   ❌ 失败: {result.stderr}")
        return
    
    print(f"   ✓ 完成")
    print("\n" + "=" * 60)
    print(f"✅ 美化成功！输出文件: {output_pptx}")
    print("=" * 60)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python test_beautify.py input.pptx")
        sys.exit(1)
    
    test_beautify(sys.argv[1])
